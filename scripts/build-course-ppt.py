from __future__ import annotations

import hashlib
import json
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "deliverables" / "DebugMate-V0.1.pptx"
MANIFEST = ROOT / "deliverables" / "asset-manifest.json"
EVIDENCE = ROOT / "evidence" / "course-v0.1" / "screenshots"

NAVY = RGBColor(18, 43, 70)
BLUE = RGBColor(42, 111, 151)
CYAN = RGBColor(43, 166, 181)
ORANGE = RGBColor(235, 113, 62)
PALE = RGBColor(239, 245, 249)
MID = RGBColor(91, 108, 122)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(49, 136, 91)
AMBER = RGBColor(204, 132, 30)


def add_text(slide, x, y, w, h, text, size=20, color=NAVY, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Microsoft YaHei"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    if align is not None:
        paragraph.alignment = align
    return box


def add_title(slide, title, subtitle=None, page=None):
    add_text(slide, 0.65, 0.3, 11.7, 0.55, title, 26, NAVY, True)
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(0.91), Inches(1.0), Inches(0.06)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = ORANGE
    slide.shapes[-1].line.fill.background()
    if subtitle:
        add_text(slide, 1.82, 0.78, 10.4, 0.25, subtitle, 10, MID)
    if page is not None:
        add_text(slide, 12.15, 7.05, 0.45, 0.2, str(page), 9, MID, align=PP_ALIGN.RIGHT)


def add_bullets(slide, items, x, y, w, h, size=18, color=NAVY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(10)
        paragraph.text = f"• {item}"
    return box


def add_card(slide, x, y, w, h, title, body, accent=BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(207, 220, 229)
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = accent
    slide.shapes[-1].line.fill.background()
    add_text(slide, x + 0.25, y + 0.12, w - 0.4, 0.4, title, 16, accent, True)
    add_text(slide, x + 0.25, y + 0.55, w - 0.45, h - 0.68, body, 12, NAVY)


def add_top_crop(slide, path: Path, x, y, w, h, crop_ratio=0.48):
    with Image.open(path) as image:
        crop_height = max(1, int(image.height * crop_ratio))
        cropped = image.crop((0, 0, image.width, crop_height))
        stream = BytesIO()
        cropped.save(stream, format="PNG")
        stream.seek(0)
        ratio = cropped.width / cropped.height
        box_ratio = w / h
        if ratio > box_ratio:
            pic_w = w
            pic_h = w / ratio
            pic_x = x
            pic_y = y + (h - pic_h) / 2
        else:
            pic_h = h
            pic_w = h * ratio
            pic_y = y
            pic_x = x + (w - pic_w) / 2
        slide.shapes.add_picture(stream, Inches(pic_x), Inches(pic_y), Inches(pic_w), Inches(pic_h))


def base_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = PALE
    return slide


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = base_slide(prs)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.0), Inches(0.14), Inches(5.5)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = ORANGE
    slide.shapes[-1].line.fill.background()
    add_text(slide, 1.18, 1.45, 10.8, 1.2, "DebugMate", 46, WHITE, True)
    add_text(
        slide, 1.2, 2.55, 10.8, 0.9, "面向 AI 专业学习场景的多模态报错诊断与复盘智能体", 25, WHITE
    )
    add_text(
        slide,
        1.2,
        4.0,
        9.4,
        0.55,
        "《校外实训》课程项目 · V0.1 本地演示版",
        18,
        RGBColor(190, 222, 235),
        True,
    )
    add_text(
        slide,
        1.2,
        5.45,
        10.5,
        0.5,
        "有依据 · 可执行 · 说明不确定性 · 文字 / PNG / MP3 同源",
        15,
        WHITE,
    )

    slide = base_slide(prs)
    add_title(slide, "为什么要做 DebugMate？", "问题背景", 2)
    add_card(
        slide,
        0.8,
        1.35,
        3.75,
        4.7,
        "报错信息复杂",
        "Traceback、环境、版本、路径和截图分散，学生很难判断应该先检查哪一项。",
        BLUE,
    )
    add_card(
        slide,
        4.78,
        1.35,
        3.75,
        4.7,
        "通用回答缺证据",
        "普通对话容易给出看似合理的安装命令，却没有绑定当前环境和官方资料。",
        ORANGE,
    )
    add_card(
        slide,
        8.76,
        1.35,
        3.75,
        4.7,
        "复盘难以沉淀",
        "一次修好不等于真正理解；缺少结构化报告、流程图和语音总结。",
        CYAN,
    )
    add_text(
        slide,
        1.0,
        6.35,
        11.2,
        0.45,
        "目标：把一次报错转化为可追溯、可检查、可复盘的学习过程。",
        20,
        NAVY,
        True,
        PP_ALIGN.CENTER,
    )

    slide = base_slide(prs)
    add_title(slide, "V0.1 的目标与边界", "按课程价值收束，而非生产发布", 3)
    add_card(
        slide,
        0.85,
        1.35,
        5.65,
        4.8,
        "必须完成",
        (
            "• Windows 本地浏览器真实演示\n• 专属知识库与完整工作流\n"
            "• 上传前自动脱敏 + 预览确认\n• 同源文字报告、PNG、MP3、ZIP\n"
            "• 3–5 个代表性案例\n• PPT 与 3 分钟以上讲解视频"
        ),
        GREEN,
    )
    add_card(
        slide,
        6.82,
        1.35,
        5.65,
        4.8,
        "明确延后",
        (
            "• 公网部署、账户与权限\n• SLA、监控、并发和跨平台\n"
            "• 完整 15 行视觉认证\n• 原子 evidence generation/pointer\n"
            "• 故障注入和攻击矩阵\n• 必须依赖付费云端 API"
        ),
        AMBER,
    )

    slide = base_slide(prs)
    add_title(slide, "工具分工与总体架构", "云端可选增强，本地闭环保证可演示", 4)
    labels = [
        ("多模态输入", "文本 / 截图 / 代码 / 环境", BLUE),
        ("隐私预览", "OCR + 规则脱敏 + 用户确认", CYAN),
        ("诊断工作流", "分类 + 知识检索 + 结构化推理", ORANGE),
        ("单一事实源", "DiagnosisRecord 1.1", BLUE),
        ("三模态结果", "报告 + PNG + MP3 + ZIP", GREEN),
    ]
    x = 0.55
    for index, (title, body, color) in enumerate(labels):
        add_card(slide, x, 2.05, 2.25, 2.45, title, body, color)
        if index < len(labels) - 1:
            add_text(slide, x + 2.28, 2.9, 0.35, 0.4, "→", 24, ORANGE, True, PP_ALIGN.CENTER)
        x += 2.55
    add_text(
        slide,
        1.15,
        5.25,
        11.0,
        0.85,
        (
            "Dify Cloud：视觉 / RAG / LLM 编排（可选）    |    "
            "本地 Python：脱敏 / 校验 / 产物 / Gradio（稳定主闭环）"
        ),
        16,
        NAVY,
        True,
        PP_ALIGN.CENTER,
    )

    slide = base_slide(prs)
    add_title(slide, "专属知识库：只用可核验官方来源", "17 个来源，覆盖 7 类技术主题", 5)
    topics = [
        "Python 异常与导入",
        "pip / venv",
        "PyTorch 张量",
        "CUDA 与显存",
        "Hugging Face",
        "Ultralytics",
        "Windows 路径 / PowerShell",
    ]
    for i, topic in enumerate(topics):
        row, col = divmod(i, 4)
        add_card(
            slide,
            0.75 + col * 3.12,
            1.4 + row * 2.1,
            2.82,
            1.65,
            topic,
            "保存标题、URL、版本范围、错误类别与选择理由",
            [BLUE, CYAN, ORANGE, GREEN][col],
        )
    add_text(
        slide,
        0.95,
        5.85,
        11.4,
        0.65,
        "诊断输出引用命中的 evidence_id；没有证据时写明缺失信息，而不是补造确定结论。",
        17,
        NAVY,
        True,
        PP_ALIGN.CENTER,
    )

    slide = base_slide(prs)
    add_title(slide, "隐私安全与完整工作流", "自动脱敏 + 上传前预览确认", 6)
    steps = [
        "本机校验输入",
        "OCR / 正则发现敏感项",
        "生成脱敏文本与遮挡图",
        "用户预览并确认",
        "事实抽取与纠错",
        "知识检索与诊断",
        "多模态派生与二次扫描",
    ]
    for i, step in enumerate(steps):
        y = 1.25 + i * 0.76
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.95), Inches(y), Inches(0.48), Inches(0.48)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ORANGE
        circle.line.fill.background()
        add_text(slide, 0.95, y, 0.48, 0.48, str(i + 1), 13, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, 1.65, y - 0.02, 5.0, 0.5, step, 17, NAVY, True)
    add_card(
        slide,
        7.3,
        1.4,
        5.05,
        4.7,
        "默认保护对象",
        (
            "Token / 密码\n邮箱 / 用户名\nWindows 与 Linux 绝对路径\n"
            "截图中的敏感像素区域\n日志与导出物中的二次泄漏"
        ),
        CYAN,
    )

    slide = base_slide(prs)
    add_title(slide, "提示词 V1–V4：从“能回答”到“可派生”", "优化重点是结构、证据和一致性", 7)
    versions = [
        ("V1", "基础诊断", "类别 / 原因 / 步骤\n问题：格式与依据不稳定", MID),
        ("V2", "引用约束", "fact/evidence 绑定\n事实 / 推断 / 缺失信息分离", BLUE),
        ("V3", "结构可靠", "DiagnosisRecord JSON\n命令影响 / 预期 / 回退", ORANGE),
        ("V4", "课程定稿", "recap 长度与中英分工\n报告 / PNG / MP3 同源", GREEN),
    ]
    for i, (v, title, body, color) in enumerate(versions):
        add_card(slide, 0.65 + i * 3.13, 1.55, 2.83, 3.8, f"{v} · {title}", body, color)
    add_text(
        slide,
        0.9,
        5.75,
        11.6,
        0.7,
        (
            "真实性说明：V1 已由固定案例 manifest 记录；"
            "V2–V4 为可直接配置的设计迭代，不伪造云端批量分数。"
        ),
        15,
        NAVY,
        True,
        PP_ALIGN.CENTER,
    )

    slide = base_slide(prs)
    add_title(slide, "完整成果：统一诊断工作台", "当前代码真实 Microsoft Edge 截图", 8)
    add_top_crop(slide, EVIDENCE / "01-completed-overview.png", 0.65, 1.15, 8.15, 5.75, 0.58)
    add_card(
        slide,
        9.05,
        1.35,
        3.65,
        4.8,
        "页面同时展示",
        (
            "• 脱敏输入与抽取字段\n• 错误类别与置信度\n• 事实和官方引用\n"
            "• 文字报告 / 诊断卡 / 语音\n• 来源运行标识\n• 完整 ZIP 下载"
        ),
        BLUE,
    )
    add_text(
        slide, 9.1, 6.35, 3.55, 0.4, "页面明确标注“离线回放”", 14, ORANGE, True, PP_ALIGN.CENTER
    )

    slide = base_slide(prs)
    add_title(slide, "三模态由同一诊断对象派生", "不是装饰图片或无关配音", 9)
    add_card(
        slide,
        0.8,
        1.45,
        3.75,
        4.75,
        "文字报告",
        "事实、根因候选、缺失信息、检查、修复、验证、置信度与局限。",
        BLUE,
    )
    add_card(
        slide,
        4.8,
        1.45,
        3.75,
        4.75,
        "PNG 诊断卡",
        "Pillow 确定性绘制，保留类别、证据关系和关键步骤；不让图像模型改写证据。",
        ORANGE,
    )
    add_card(
        slide,
        8.8,
        1.45,
        3.75,
        4.75,
        "MP3 语音复盘",
        "从同一 recap_text 合成；记录后端、时长和降级原因，并用 FFmpeg 验证。",
        GREEN,
    )
    add_text(
        slide,
        1.05,
        6.4,
        11.2,
        0.4,
        "共同身份：case_id + source_run_id + diagnosis hash + schema version",
        17,
        NAVY,
        True,
        PP_ALIGN.CENTER,
    )

    slide = base_slide(prs)
    add_title(slide, "失败也要诚实：部分完成与最小重试", "当前代码真实 Edge 截图", 10)
    add_top_crop(slide, EVIDENCE / "02-tts-partial.png", 0.55, 1.2, 5.95, 5.4, 0.52)
    add_top_crop(slide, EVIDENCE / "03-card-partial.png", 6.83, 1.2, 5.95, 5.4, 0.52)
    add_text(
        slide,
        0.75,
        6.55,
        5.55,
        0.35,
        "TTS 失败：保留报告与诊断卡，只重试语音",
        13,
        AMBER,
        True,
        PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        7.05,
        6.55,
        5.5,
        0.35,
        "PNG 失败：保留报告与语音，只重试诊断卡",
        13,
        AMBER,
        True,
        PP_ALIGN.CENTER,
    )

    slide = base_slide(prs)
    add_title(slide, "V0.1 代表性验证结果", "课程演示所需，而非生产发布认证", 11)
    metrics = [
        ("58", "界面 / 状态 / 回调测试", GREEN),
        ("264", "报告 / PNG / MP3 / 结果模块测试", BLUE),
        ("4+", "代表性真实 Edge 场景", ORANGE),
        ("1", "真实 ZIP 同次运行验证", CYAN),
    ]
    for i, (number, label, color) in enumerate(metrics):
        x = 0.72 + i * 3.13
        add_card(slide, x, 1.45, 2.82, 2.15, number, label, color)
        add_text(slide, x + 0.15, 1.72, 2.5, 0.75, number, 36, color, True, PP_ALIGN.CENTER)
    add_bullets(
        slide,
        [
            "键盘顺序与状态播报",
            "200% 缩放与无整体横向溢出",
            "长报告、长命令和高诊断卡",
            "成功 / 部分完成 / 失败状态辨识",
            "ZIP 文件名、MIME、manifest、校验值与 source_run_id",
        ],
        1.15,
        4.05,
        11.0,
        2.2,
        17,
    )

    slide = base_slide(prs)
    add_title(slide, "局限与后续改进", "如实说明大模型工具的边界", 12)
    add_card(
        slide,
        0.75,
        1.35,
        3.75,
        4.85,
        "当前局限",
        (
            "• Dify 在线能力受账号与额度影响\n• 案例数量少，不能代表全面覆盖\n"
            "• SAPI 正式录制前需人工试听\n• 回放不等同于云端实时推理"
        ),
        AMBER,
    )
    add_card(
        slide,
        4.8,
        1.35,
        3.75,
        4.85,
        "近期改进",
        (
            "• 完成 Dify 图像 + RAG smoke test\n• 增加 CUDA 显存和张量形状案例\n"
            "• 用同案例对照 V1–V4\n• 继续精简首屏和录屏操作"
        ),
        BLUE,
    )
    add_card(
        slide,
        8.85,
        1.35,
        3.75,
        4.85,
        "长期方向",
        (
            "• 更多框架与错误类别\n• 个性化学习复盘\n"
            "• 可信度校准与人机协作\n• 经批准后再考虑部署与监控"
        ),
        GREEN,
    )

    slide = base_slide(prs)
    add_title(slide, "总结", "DebugMate V0.1 已形成完整课程闭环", 13)
    add_text(
        slide,
        1.0,
        1.5,
        11.3,
        0.8,
        "一次真实或可复现的 AI/Python 报错",
        26,
        NAVY,
        True,
        PP_ALIGN.CENTER,
    )
    add_text(slide, 1.0, 2.5, 11.3, 0.55, "↓", 30, ORANGE, True, PP_ALIGN.CENTER)
    add_text(
        slide,
        1.0,
        3.1,
        11.3,
        0.85,
        "隐私确认 + 官方知识 + 结构化诊断",
        27,
        BLUE,
        True,
        PP_ALIGN.CENTER,
    )
    add_text(slide, 1.0, 4.05, 11.3, 0.55, "↓", 30, ORANGE, True, PP_ALIGN.CENTER)
    add_text(
        slide,
        1.0,
        4.65,
        11.3,
        0.85,
        "文字报告 + PNG 诊断卡 + MP3 语音复盘 + ZIP",
        25,
        GREEN,
        True,
        PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        1.0,
        6.1,
        11.3,
        0.5,
        "帮助 AI 专业学生把“修好一次”变成“理解并复盘一次”",
        19,
        NAVY,
        True,
        PP_ALIGN.CENTER,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    ppt_hash = hashlib.sha256(OUT.read_bytes()).hexdigest()
    assets = [
        ROOT / "docs" / "course" / "README.md",
        ROOT / "docs" / "course" / "presentation-outline.md",
        ROOT / "docs" / "course" / "video-script.md",
        ROOT / "docs" / "course" / "prompt-iteration.md",
        *sorted(EVIDENCE.glob("*.png")),
        OUT,
    ]
    manifest = {
        "schema_version": "debugmate-course-assets-1.0",
        "generated_on": "2026-07-19",
        "pptx": {
            "path": str(OUT.relative_to(ROOT)).replace("\\", "/"),
            "sha256": ppt_hash,
            "slides": len(prs.slides),
        },
        "assets": [
            {
                "path": str(path.relative_to(ROOT)).replace("\\", "/"),
                "bytes": path.stat().st_size,
                "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
            }
            for path in assets
        ],
    }
    MANIFEST.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


if __name__ == "__main__":
    build()
